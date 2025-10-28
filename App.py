# app.py
import os
import tempfile
import datetime
from io import BytesIO

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.figure_factory as ff
import streamlit as st
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
from sklearn.linear_model import LinearRegression

# ----------------------------
# PAGE CONFIG
# ----------------------------
st.set_page_config(page_title="Auto-Dashboard Generator", layout="wide", page_icon="📊")
st.title("📊 Auto Dashboard Generator – With AI Insights")

st.markdown("""
<style>
    [data-testid="stHeader"] {
        background: linear-gradient(90deg, #3a0ca3, #4361ee);
    }
    [data-testid="stSidebar"] {
        background-color: #1a1a2e;
    }
    [data-testid="stMarkdownContainer"] {
        color: #f0f0f0;
    }
</style>
""", unsafe_allow_html=True)

st.markdown("Upload → Clean → Auto & Custom Charts → Export Dashboard (PDF/PPT) with AI Insights")

# ----------------------------
# UTILITY FUNCTIONS
# ----------------------------
def safe_color_sequence(theme_name):
    try:
        if theme_name == "Plotly":
            return px.colors.qualitative.Plotly
        return getattr(px.colors.sequential, theme_name)
    except Exception:
        return px.colors.qualitative.Plotly

def save_temp_fig_png(fig):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
    tmp.close()
    fig.write_image(tmp.name, format='png', engine='kaleido')
    return tmp.name

def compute_trend_over_time(df, date_col, value_col):
    tmp = df[[date_col, value_col]].dropna().copy()
    if tmp.empty:
        return None
    tmp[date_col] = pd.to_datetime(tmp[date_col], errors='coerce')
    tmp = tmp.dropna()
    if tmp.empty:
        return None
    monthly = tmp.set_index(date_col).resample('M').mean().reset_index()
    if len(monthly) < 3:
        return None
    X = np.arange(len(monthly)).reshape(-1, 1)
    y = monthly[value_col].values.reshape(-1, 1)
    model = LinearRegression().fit(X, y)
    slope = float(model.coef_[0])
    if slope > 0:
        return f"Upward trend detected for {value_col}."
    elif slope < 0:
        return f"Downward trend detected for {value_col}."
    else:
        return f"No significant trend detected for {value_col}."

def generate_summary_insights(df):
    if df is None:
        return []
    insights = []
    numeric_cols = df.select_dtypes(include=['float64','int64']).columns.tolist()
    categorical_cols = df.select_dtypes(include=['object','category']).columns.tolist()

    for col in numeric_cols:
        mean_val = df[col].mean()
        max_val = df[col].max()
        min_val = df[col].min()
        insights.append(f"{col}: mean={mean_val:.2f}, max={max_val}, min={min_val}")

    for col in categorical_cols:
        top = df[col].mode()[0] if not df[col].mode().empty else 'N/A'
        count = df[col].value_counts().iloc[0] if not df[col].value_counts().empty else 0
        insights.append(f"{col}: most frequent = {top} ({count} times)")

    datetime_cols = df.select_dtypes(include=['datetime64']).columns.tolist()
    for date_col in datetime_cols:
        for num_col in numeric_cols:
            trend = compute_trend_over_time(df, date_col, num_col)
            if trend:
                insights.append(trend)

    return insights

# ----------------------------
# FILE UPLOAD
# ----------------------------
if "df" not in st.session_state:
    st.session_state.df = None
if "original_df" not in st.session_state:
    st.session_state.original_df = None

uploaded_file = st.file_uploader("Upload CSV or Excel file", type=["csv", "xlsx"])

if uploaded_file is not None:
    try:
        if uploaded_file.name.endswith('.csv'):
            st.session_state.df = pd.read_csv(uploaded_file)
        else:
            xls = pd.ExcelFile(uploaded_file)
            if len(xls.sheet_names) == 0:
                st.error("⚠️ Uploaded Excel file has no sheets.")
                st.stop()
            st.session_state.df = pd.read_excel(xls, sheet_name=xls.sheet_names[0])
        st.session_state.original_df = st.session_state.df.copy()
        st.success(f"✅ Uploaded: {uploaded_file.name} — {st.session_state.df.shape[0]} rows × {st.session_state.df.shape[1]} cols")
    except Exception as e:
        st.error(f"⚠️ Error reading file: {e}")
        st.stop()
elif st.session_state.df is not None:
    st.info("📁 Using previously uploaded data.")
else:
    st.info("👈 Please upload a CSV or Excel file to start.")
    st.stop()

df = st.session_state.df
original_df = st.session_state.original_df
insights_list = generate_summary_insights(df)

# ----------------------------
# TABS
# ----------------------------
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "🧹 Cleaning",
    "📈 Auto Charts",
    "🧠 Smart Insights",
    "⚙️ Custom Charts",
    "💾 Export Dashboard"
])

chart_list = []

# ----------------------------
# TAB 1: CLEANING
# ----------------------------
with tab1:
    st.subheader("Data Cleaning")
    with st.expander("👀 Raw Data (First 10 Rows)"):
        st.dataframe(df.head(10), use_container_width=True)

    if st.checkbox("Remove Duplicates"):
        before = len(df)
        df = df.drop_duplicates()
        st.success(f"Removed {before - len(df)} duplicates")

    missing_action = st.selectbox("Handle Missing Values", ["Do Nothing","Drop Rows","Fill Mean","Fill Median","Fill Mode"])
    if missing_action != "Do Nothing":
        numeric_cols = df.select_dtypes(include=['float64','int64']).columns
        if missing_action=="Drop Rows": df=df.dropna()
        elif missing_action=="Fill Mean": df[numeric_cols]=df[numeric_cols].fillna(df[numeric_cols].mean())
        elif missing_action=="Fill Median": df[numeric_cols]=df[numeric_cols].fillna(df[numeric_cols].median())
        elif missing_action=="Fill Mode":
            for col in df.columns: df[col].fillna(df[col].mode()[0], inplace=True)
        st.success("✅ Missing values handled")

    rename_col = st.selectbox("Select column to rename", options=["None"] + list(df.columns))
    if rename_col != "None":
        new_name = st.text_input(f"Rename '{rename_col}' to:")
        if new_name:
            df.rename(columns={rename_col:new_name}, inplace=True)
            st.success(f"Renamed {rename_col} → {new_name}")

    col_to_convert = st.selectbox("Select column to convert type", options=["None"] + list(df.columns))
    dtype = st.selectbox("Convert to:", ["int","float","string"])
    if st.button("Convert Type") and col_to_convert != "None":
        try:
            if dtype in ["int","float"]:
                df[col_to_convert] = pd.to_numeric(df[col_to_convert], errors='coerce')
                if dtype=="int": df[col_to_convert]=df[col_to_convert].fillna(0).astype(int)
                else: df[col_to_convert]=df[col_to_convert].astype(float)
            elif dtype=="string": df[col_to_convert]=df[col_to_convert].astype(str)
            st.success(f"✅ Converted '{col_to_convert}' to {dtype}")
        except Exception as e:
            st.error(f"⚠️ Error: {e}")

    csv = df.to_csv(index=False).encode('utf-8')
    st.download_button("💾 Download Cleaned CSV", data=csv, file_name="cleaned_data.csv")

# ----------------------------
# TAB 2: AUTO CHARTS
# ----------------------------
with tab2:
    st.subheader("Auto Charts")
    color_theme = st.selectbox("Color Theme", ["Plotly","Viridis","Cividis","Plasma","Inferno"], key="auto_color")

    numeric_cols = df.select_dtypes(include=['int64','float64']).columns.tolist()
    categorical_cols = df.select_dtypes(include=['object','category']).columns.tolist()

    for col in df.columns:
        if df[col].dtype == 'object':
            try:
                df[col] = pd.to_datetime(df[col])
            except Exception:
                pass
    datetime_cols = df.select_dtypes(include=['datetime64[ns]','datetime64']).columns.tolist()

    st.subheader("📈 Key Performance Indicators")
    kpicol1, kpicol2, kpicol3 = st.columns(3)
    if numeric_cols:
        kpicol1.metric("Total", f"{df[numeric_cols[0]].sum():,.2f}")
        kpicol2.metric("Average", f"{df[numeric_cols[0]].mean():,.2f}")
        kpicol3.metric("Count", f"{len(df):,}")

    # Auto Charts
    if numeric_cols and categorical_cols:
        color_seq = safe_color_sequence(color_theme)
        fig = px.bar(df, x=categorical_cols[0], y=numeric_cols[0], title=f"{numeric_cols[0]} by {categorical_cols[0]}", color_discrete_sequence=color_seq)
        st.plotly_chart(fig, use_container_width=True)
        chart_list.append((fig, f"{numeric_cols[0]} by {categorical_cols[0]}"))

    if numeric_cols and datetime_cols:
        fig = px.line(df.sort_values(datetime_cols[0]), x=datetime_cols[0], y=numeric_cols[0], title=f"{numeric_cols[0]} over time")
        st.plotly_chart(fig, use_container_width=True)
        chart_list.append((fig, f"{numeric_cols[0]} over time"))

    if numeric_cols:
        fig = px.histogram(df, x=numeric_cols[0], title=f"Distribution of {numeric_cols[0]}")
        st.plotly_chart(fig, use_container_width=True)
        chart_list.append((fig, f"Distribution of {numeric_cols[0]}"))

    if categorical_cols:
        fig = px.pie(df, names=categorical_cols[0], title=f"{categorical_cols[0]} share")
        st.plotly_chart(fig, use_container_width=True)
        chart_list.append((fig, f"{categorical_cols[0]} share"))

# ----------------------------
# TAB 3: SMART INSIGHTS
# ----------------------------
with tab3:
    st.title("🧠 Smart Insights Dashboard")

    if df is None or df.empty:
        st.warning("⚠️ Please upload a dataset first.")
    else:
        st.subheader("📊 Dataset Overview")
        st.dataframe(df.describe().T.style.highlight_max(axis=0), use_container_width=True)

        st.divider()
        st.subheader("🔍 Correlation Heatmap")

        numeric_cols = df.select_dtypes(include=['float64', 'int64']).columns
        if len(numeric_cols) > 1:
            corr = df[numeric_cols].corr()
            fig_corr = ff.create_annotated_heatmap(
                z=corr.values,
                x=list(corr.columns),
                y=list(corr.index),
                colorscale='Viridis',
                showscale=True
            )
            st.plotly_chart(fig_corr, use_container_width=True)
        else:
            st.info("ℹ️ Not enough numeric columns to generate a correlation heatmap.")

        st.divider()
        st.subheader("📈 Trend Insights")

        if len(numeric_cols) >= 2:
            diffs = df[numeric_cols].iloc[-1] - df[numeric_cols].iloc[0]
            top_growth = diffs.idxmax()
            top_decline = diffs.idxmin()

            col1, col2 = st.columns(2)
            col1.success(f"📈 **{top_growth}** shows the highest growth across the dataset.")
            col2.error(f"📉 **{top_decline}** shows a decline over time.")
        else:
            st.info("ℹ️ Need at least two numeric columns to detect trends.")

# ----------------------------
# TAB 4: CUSTOM CHARTS
# ----------------------------
with tab4:
    st.subheader("Custom Charts")
    chart_type = st.selectbox("Chart Type", ["Bar","Line","Scatter","Histogram","Pie"], key="custom_type")
    x_col = st.selectbox("X-axis", df.columns, key="custom_x")
    y_col = st.selectbox("Y-axis (optional)", ["None"] + list(df.columns), key="custom_y")
    color_theme_custom = st.selectbox("Color Theme", ["Plotly","Viridis","Cividis","Plasma","Inferno"], key="custom_color")

    if st.button("Generate Chart", key="custom_chart_btn"):
        try:
            color_seq = safe_color_sequence(color_theme_custom)
            if chart_type == "Bar":
                fig = px.bar(df, x=x_col, y=(None if y_col == 'None' else y_col), color_discrete_sequence=color_seq)
            elif chart_type == "Line":
                fig = px.line(df, x=x_col, y=(None if y_col == 'None' else y_col))
            elif chart_type == "Scatter":
                fig = px.scatter(df, x=x_col, y=(None if y_col == 'None' else y_col))
            elif chart_type == "Histogram":
                fig = px.histogram(df, x=x_col)
            elif chart_type == "Pie":
                fig = px.pie(df, names=x_col)
            st.plotly_chart(fig, use_container_width=True)
            chart_list.append((fig, f"Custom: {chart_type} - {x_col} {y_col}"))
        except Exception as e:
            st.error(f"⚠️ Error: {e}")

# ----------------------------
# TAB 5: EXPORT DASHBOARD
# ----------------------------
with tab5:
    st.subheader("Export Dashboard with AI Insights")

    if st.button("Export as PDF"):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial","B",16)
        pdf.cell(0,10,"Dashboard Export",0,1,'C')
        pdf.ln(5)

        pdf.set_font("Arial","",12)
        for insight in insights_list:
            pdf.multi_cell(0,8,insight)
        pdf.ln(5)

        for fig, title in chart_list:
            img_file = save_temp_fig_png(fig)
            pdf.image(img_file, w=90)
            os.remove(img_file)

        pdf_output = BytesIO()
        pdf.output(pdf_output)
        pdf_output.seek(0)
        st.download_button("💾 Download PDF", data=pdf_output, file_name="dashboard.pdf", mime="application/pdf")

    if st.button("Export as PPTX"):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]

        for insight in insights_list:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Insight"
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            textbox.text = insight

        for fig, title in chart_list:
            slide = prs.slides.add_slide(slide_layout)
            img_file = save_temp_fig_png(fig)
            slide.shapes.add_picture(img_file, Inches(1), Inches(1.25), width=Inches(8))
            os.remove(img_file)

        pptx_output = BytesIO()
        prs.save(pptx_output)
        pptx_output.seek(0)
        st.download_button("💾 Download PPTX", data=pptx_output, file_name="dashboard.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
